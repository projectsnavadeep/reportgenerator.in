import telebot
import threading
import time
import os
from telebot import types
from report_generator import ReportGenerator

from docx import Document
from pptx import Presentation
from pptx.util import Pt
from reportlab.lib.pagesizes import A4
from reportlab.platypus import SimpleDocTemplate, Paragraph
from reportlab.lib.styles import getSampleStyleSheet

BOT_TOKEN = "8455804390:AAHvaU0MRRVPAyJHnxR4yqZnZ9X3G1ASxLo"

bot = telebot.TeleBot(BOT_TOKEN)
generator = ReportGenerator()

# ---------------- SESSION MEMORY ----------------
users = {}  # chat_id -> dict


def get_user(chat_id):
    if chat_id not in users:
        users[chat_id] = {}
    return users[chat_id]


# ---------------- FILE GENERATORS ----------------
def create_txt(report_text, filename):
    with open(filename, "w", encoding="utf-8") as f:
        f.write(report_text)


def create_docx(report_text, filename):
    doc = Document()
    for line in report_text.split("\n"):
        doc.add_paragraph(line)
    doc.save(filename)


def create_ppt(report_text, filename):
    prs = Presentation()
    layout = prs.slide_layouts[1]

    for block in report_text.split("\n\n"):
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = "Business Report"
        body = slide.placeholders[1].text_frame
        body.clear()

        for line in block.split("\n"):
            p = body.add_paragraph()
            p.text = line
            p.font.size = Pt(14)

    prs.save(filename)


def create_pdf(report_text, filename):
    styles = getSampleStyleSheet()
    doc = SimpleDocTemplate(filename, pagesize=A4)
    story = []
    for line in report_text.split("\n"):
        story.append(Paragraph(line.replace("&", "&amp;"), styles["Normal"]))
    doc.build(story)


# ---------------- ANIMATION ----------------
def typing(chat_id, stop_event):
    while not stop_event.is_set():
        bot.send_chat_action(chat_id, "typing")
        time.sleep(3)


# ---------------- REPORT GENERATION ----------------
def generate_report_async(chat_id):
    user = get_user(chat_id)
    stop_event = threading.Event()
    threading.Thread(target=typing, args=(chat_id, stop_event)).start()

    try:
        report = generator.generate_report(
            data_content=user["content"],
            data_type=user["data_type"],
            user_name="User",
            user_role="Owner",
            business_type="General Business",
            report_focus="FULL"
        )

        stop_event.set()

        fmt = user.get("output_format", "txt")
        base = f"Business_Report_{chat_id}"
        file_map = {
            "txt": (f"{base}.txt", create_txt),
            "docx": (f"{base}.docx", create_docx),
            "ppt": (f"{base}.pptx", create_ppt),
            "pdf": (f"{base}.pdf", create_pdf),
        }

        filename, creator = file_map[fmt]
        creator(report, filename)

        bot.send_message(chat_id, f"✅ Report ready ({fmt.upper()})")
        bot.send_document(chat_id, open(filename, "rb"))

        os.remove(filename)

    except Exception as e:
        stop_event.set()
        bot.send_message(chat_id, f"❌ Error: {e}")


# ---------------- FORMAT SELECTION ----------------
def ask_output_format(chat_id):
    kb = types.InlineKeyboardMarkup(row_width=2)
    kb.add(
        types.InlineKeyboardButton("📄 PDF", callback_data="fmt_pdf"),
        types.InlineKeyboardButton("📝 Word", callback_data="fmt_docx"),
        types.InlineKeyboardButton("📊 PPT", callback_data="fmt_ppt"),
        types.InlineKeyboardButton("📃 Text", callback_data="fmt_txt"),
    )
    bot.send_message(chat_id, "Choose output format:", reply_markup=kb)


@bot.callback_query_handler(func=lambda c: c.data.startswith("fmt_"))
def choose_format(call):
    fmt = call.data.replace("fmt_", "")
    get_user(call.message.chat.id)["output_format"] = fmt
    bot.send_message(call.message.chat.id, "🧠 Generating report…")
    threading.Thread(
        target=generate_report_async,
        args=(call.message.chat.id,)
    ).start()


# ---------------- START ----------------
@bot.message_handler(commands=["start"])
def start(message):
    bot.reply_to(
        message,
        "📊 Send CSV / Excel / JSON / Text.\n"
        "Then choose output format (PDF / Word / PPT / Text).\n\n"
        "No questions. Direct results.",
    )


# ---------------- INPUT HANDLING ----------------
@bot.message_handler(content_types=["text"])
def receive_text(message):
    user = get_user(message.chat.id)
    user["content"] = message.text
    user["data_type"] = "text"
    ask_output_format(message.chat.id)


@bot.message_handler(content_types=["document"])
def receive_file(message):
    user = get_user(message.chat.id)

    file_bytes = bot.download_file(
        bot.get_file(message.document.file_id).file_path
    )
    filename = message.document.file_name.lower()

    if filename.endswith(".csv"):
        user["data_type"] = "csv"
        user["content"] = file_bytes.decode("utf-8", errors="ignore")
    elif filename.endswith(".json"):
        user["data_type"] = "json"
        user["content"] = file_bytes.decode("utf-8", errors="ignore")
    elif filename.endswith(".xlsx"):
        user["data_type"] = "excel"
        user["content"] = file_bytes
    else:
        bot.send_message(message.chat.id, "❌ Unsupported file type")
        return

    ask_output_format(message.chat.id)


# ---------------- RUN ----------------
if __name__ == "__main__":
    print("🤖 Brutal Telegram Bot running (PDF / DOCX / PPT / TXT)")
    bot.infinity_polling(skip_pending=True)
