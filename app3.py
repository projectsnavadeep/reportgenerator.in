"""
AI Business Companion - Professional Edition
Perplexity AI-inspired minimal design with intelligent data processing
Optimized for all devices, production-ready
"""

from flask import Flask, request, jsonify, send_file, session, redirect, url_for
from report_generator import ReportGenerator
from auth import init_db, create_user, validate_user
import os
import io
import re
import json
import pandas as pd
import tempfile
import shutil

from datetime import datetime

# Document processing
from reportlab.lib.pagesizes import A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_CENTER, TA_LEFT
from reportlab.lib import colors
from docx import Document
from docx.shared import Inches, RGBColor, Pt
from pptx import Presentation
from pptx.util import Pt as PptxPt
from pptx.enum.text import PP_ALIGN

app = Flask(__name__)
app.secret_key = os.environ.get('SECRET_KEY') or 'dev-secret-key-change-in-production'
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 50MB max file size
init_db()
generator = ReportGenerator()

# ---------------- INTELLIGENT FILE PROCESSOR ----------------

class DataProcessor:
    """Universal data processor that auto-detects format and converts to usable text"""
    
    SUPPORTED_EXTENSIONS = {
        'csv', 'xlsx', 'xls', 'json', 'txt', 'md', 'pdf', 'docx', 'doc'
    }
    
    @staticmethod
    def detect_format(content: bytes, filename: str = "") -> str:
        """Auto-detect data format from content"""
        if not content:
            return "empty"
        
        # Try to decode as text first
        try:
            text = content.decode('utf-8', errors='ignore').strip()
            
            # Check for JSON
            if text.startswith('{') or text.startswith('['):
                try:
                    json.loads(text)
                    return "json"
                except:
                    pass
            
            # Check for CSV
            lines = text.split('\n')
            if len(lines) > 1:
                # Check for comma-separated or tab-separated
                first_line = lines[0]
                if ',' in first_line or '\t' in first_line:
                    return "csv"
            
            # Check for markdown/table format
            if '|' in text and '\n' in text:
                return "markdown"
            
            return "text"
        except:
            pass
        
        # Check filename extension
        ext = filename.lower().split('.')[-1] if '.' in filename else ''
        if ext in ['xlsx', 'xls']:
            return "excel"
        elif ext == 'pdf':
            return "pdf"
        elif ext in ['docx', 'doc']:
            return "word"
        
        return "binary"
    
    @staticmethod
    def extract_text(content: bytes, filename: str = "", detected_format: str = None) -> str:
        """Extract readable text from any file format"""
        fmt = detected_format or DataProcessor.detect_format(content, filename)
        
        if fmt == "empty":
            raise ValueError("File is empty")
        
        if fmt == "json":
            try:
                data = json.loads(content.decode('utf-8', errors='ignore'))
                return DataProcessor._flatten_json(data)
            except Exception as e:
                return content.decode('utf-8', errors='ignore')
        
        elif fmt == "csv":
            return content.decode('utf-8', errors='ignore')
        
        elif fmt == "excel" or fmt == "xlsx":
            try:
                import pandas as pd
                df = pd.read_excel(io.BytesIO(content))
                return df.to_csv(index=False)
            except Exception as e:
                raise ValueError(f"Could not read Excel file: {str(e)}")
        
        elif fmt == "pdf":
            try:
                import PyPDF2
                reader = PyPDF2.PdfReader(io.BytesIO(content))
                text = ""
                for page in reader.pages:
                    text += page.extract_text() + "\n"
                return text if text.strip() else "PDF contains no extractable text"
            except Exception as e:
                raise ValueError(f"Could not read PDF: {str(e)}")
        
        elif fmt == "word" or fmt == "docx":
            try:
                from docx import Document
                doc = Document(io.BytesIO(content))
                return "\n".join([para.text for para in doc.paragraphs])
            except Exception as e:
                raise ValueError(f"Could not read Word document: {str(e)}")
        
        elif fmt == "text" or fmt == "markdown":
            return content.decode('utf-8', errors='ignore')
        
        else:
            # Try as text as fallback
            try:
                return content.decode('utf-8', errors='ignore')
            except:
                raise ValueError(f"Unsupported file format: {fmt}. Please use CSV, Excel, JSON, PDF, Word, or text files.")
    
    @staticmethod
    def _flatten_json(data, prefix="") -> str:
        """Flatten nested JSON to readable text"""
        lines = []
        if isinstance(data, dict):
            for key, value in data.items():
                new_key = f"{prefix}.{key}" if prefix else key
                if isinstance(value, (dict, list)):
                    lines.append(f"{new_key}:")
                    lines.append(DataProcessor._flatten_json(value, new_key))
                else:
                    lines.append(f"{new_key}: {value}")
        elif isinstance(data, list):
            for i, item in enumerate(data[:100]):  # Limit to 100 items
                lines.append(f"{prefix}[{i}]:")
                lines.append(DataProcessor._flatten_json(item, f"{prefix}[{i}]"))
        else:
            lines.append(str(data))
        return "\n".join(lines)


# ---------------- DOCUMENT GENERATORS ----------------

def create_pdf(report_text, filename):
    """Create professional PDF with modern styling"""
    doc = SimpleDocTemplate(
        filename, 
        pagesize=A4, 
        rightMargin=60, 
        leftMargin=60, 
        topMargin=60, 
        bottomMargin=40
    )
    
    styles = getSampleStyleSheet()
    
    # Custom styles
    title_style = ParagraphStyle(
        'Title',
        parent=styles['Heading1'],
        fontSize=28,
        textColor=colors.HexColor('#1F2937'),
        spaceAfter=8,
        alignment=TA_CENTER,
        fontName='Helvetica-Bold'
    )
    
    subtitle_style = ParagraphStyle(
        'Subtitle',
        parent=styles['Normal'],
        fontSize=12,
        textColor=colors.HexColor('#6B7280'),
        alignment=TA_CENTER,
        spaceAfter=30
    )
    
    heading_style = ParagraphStyle(
        'Heading',
        parent=styles['Heading2'],
        fontSize=16,
        textColor=colors.HexColor('#374151'),
        spaceAfter=12,
        spaceBefore=20,
        fontName='Helvetica-Bold'
    )
    
    body_style = ParagraphStyle(
        'Body',
        parent=styles['BodyText'],
        fontSize=11,
        leading=18,
        textColor=colors.HexColor('#4B5563'),
        alignment=TA_LEFT
    )
    
    story = []
    
    # Header
    story.append(Paragraph("Business Intelligence Report", title_style))
    story.append(Paragraph(f"Generated on {datetime.now().strftime('%B %d, %Y')}", subtitle_style))
    story.append(Spacer(1, 20))
    
    # Process content
    sections = report_text.split('\n\n')
    for section in sections:
        if not section.strip():
            continue
        
        lines = section.split('\n')
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            # Detect headings
            if line.isupper() and len(line) < 100:
                story.append(Paragraph(line, heading_style))
            elif line.startswith('**') and line.endswith('**'):
                clean = line.replace('**', '')
                story.append(Paragraph(f"<b>{clean}</b>", heading_style))
            elif line.startswith('#'):
                level = len(line) - len(line.lstrip('#'))
                text = line.lstrip('#').strip()
                story.append(Paragraph(text, heading_style if level <= 2 else body_style))
            else:
                # Regular paragraph with bold support
                formatted = line.replace('**', '<b>', 1).replace('**', '</b>', 1) if '**' in line else line
                story.append(Paragraph(formatted, body_style))
        
        story.append(Spacer(1, 12))
    
    # Footer
    def add_page_number(canvas, doc):
        canvas.saveState()
        canvas.setFont('Helvetica', 9)
        canvas.setFillColor(colors.HexColor('#9CA3AF'))
        canvas.drawRightString(A4[0] - 60, 40, f"Page {canvas.getPageNumber()}")
        canvas.drawString(60, 40, "AI Business Companion")
        canvas.restoreState()
    
    doc.build(story, onFirstPage=add_page_number, onLaterPages=add_page_number)

def create_docx(report_text, filename):
    """Create professional Word document"""
    doc = Document()
    
    # Title
    title = doc.add_heading('Business Intelligence Report', level=0)
    title.alignment = 1
    for run in title.runs:
        run.font.color.rgb = RGBColor(31, 41, 55)
        run.font.size = Pt(28)
        run.font.bold = True
    
    # Subtitle
    subtitle = doc.add_paragraph(f"Generated on {datetime.now().strftime('%B %d, %Y')}")
    subtitle.alignment = 1
    subtitle.runs[0].font.color.rgb = RGBColor(107, 114, 128)
    subtitle.runs[0].font.size = Pt(11)
    
    doc.add_paragraph()
    
    # Process content
    sections = report_text.split('\n\n')
    for section in sections:
        if not section.strip():
            continue
        
        lines = section.split('\n')
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            if line.isupper() and len(line) < 100:
                p = doc.add_heading(line, level=1)
                p.runs[0].font.color.rgb = RGBColor(55, 65, 81)
            elif line.startswith('**') and line.endswith('**'):
                clean = line.replace('**', '')
                p = doc.add_paragraph()
                p.add_run(clean).bold = True
                p.runs[0].font.color.rgb = RGBColor(55, 65, 81)
            else:
                doc.add_paragraph(line)
    
    doc.save(filename)

def create_ppt(report_text, filename):
    """Create professional PowerPoint"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Business Intelligence Report"
    p.font.size = PptxPt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(31, 41, 55)
    p.alignment = PP_ALIGN.CENTER
    
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.333), Inches(0.5))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = datetime.now().strftime('%B %d, %Y')
    p.font.size = PptxPt(18)
    p.font.color.rgb = RGBColor(107, 114, 128)
    p.alignment = PP_ALIGN.CENTER
    
    # Content slides
    blocks = [b for b in report_text.split('\n\n') if b.strip()][:8]
    
    for block in blocks:
        slide = prs.slides.add_slide(blank_layout)
        
        lines = block.split('\n')
        header = lines[0][:60] if lines[0] else "Analysis"
        
        # Header
        header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
        tf = header_box.text_frame
        p = tf.paragraphs[0]
        p.text = header.replace('**', '').replace('#', '')
        p.font.size = PptxPt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(31, 41, 55)
        
        # Content
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12.333), Inches(5.2))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, line in enumerate(lines[1:5]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            clean = line.replace('**', '').replace('- ', '• ').replace('#', '')
            p.text = clean[:120]
            p.font.size = PptxPt(20)
            p.font.color.rgb = RGBColor(75, 85, 99)
            p.space_after = PptxPt(12)
    
    prs.save(filename)


# ---------------- PERPLEXITY-INSPIRED MINIMAL CSS ----------------

PERPLEXITY_CSS = """
:root {
    --bg: #FFFFFF;
    --surface: #F9FAFB;
    --surface-hover: #F3F4F6;
    --border: #E5E7EB;
    --border-focus: #10B981;
    --text: #111827;
    --text-secondary: #6B7280;
    --text-muted: #9CA3AF;
    --primary: #10B981;
    --primary-hover: #059669;
    --primary-light: #D1FAE5;
    --accent: #3B82F6;
    --error: #EF4444;
    --error-bg: #FEE2E2;
    --radius-sm: 8px;
    --radius-md: 12px;
    --radius-lg: 16px;
    --shadow-sm: 0 1px 2px 0 rgba(0, 0, 0, 0.05);
    --shadow-md: 0 4px 6px -1px rgba(0, 0, 0, 0.1), 0 2px 4px -1px rgba(0, 0, 0, 0.06);
    --shadow-lg: 0 10px 15px -3px rgba(0, 0, 0, 0.1), 0 4px 6px -2px rgba(0, 0, 0, 0.05);
    --transition: all 0.2s cubic-bezier(0.4, 0, 0.2, 1);
}

* {
    margin: 0;
    padding: 0;
    box-sizing: border-box;
    -webkit-font-smoothing: antialiased;
    -moz-osx-font-smoothing: grayscale;
}

html {
    scroll-behavior: smooth;
}

body {
    font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, 'Helvetica Neue', Arial, sans-serif;
    background: var(--bg);
    color: var(--text);
    line-height: 1.6;
    min-height: 100vh;
}

/* Navigation - Clean & Minimal */
.nav {
    position: fixed;
    top: 0;
    left: 0;
    right: 0;
    height: 64px;
    background: rgba(255, 255, 255, 0.95);
    backdrop-filter: blur(10px);
    border-bottom: 1px solid var(--border);
    display: flex;
    align-items: center;
    justify-content: space-between;
    padding: 0 24px;
    z-index: 100;
}

.nav-brand {
    display: flex;
    align-items: center;
    gap: 12px;
    font-weight: 600;
    font-size: 18px;
    color: var(--text);
    text-decoration: none;
}

.logo-icon {
    width: 32px;
    height: 32px;
    background: linear-gradient(135deg, var(--primary), #34D399);
    border-radius: 8px;
    display: flex;
    align-items: center;
    justify-content: center;
    color: white;
    font-weight: 700;
    font-size: 14px;
}

.nav-actions {
    display: flex;
    gap: 12px;
    align-items: center;
}

/* Container */
.container {
    max-width: 900px;
    margin: 0 auto;
    padding: 96px 24px 48px;
}

/* Hero - Perplexity Style */
.hero {
    text-align: center;
    margin-bottom: 48px;
}

.hero h1 {
    font-size: clamp(32px, 5vw, 48px);
    font-weight: 700;
    color: var(--text);
    margin-bottom: 16px;
    letter-spacing: -0.02em;
    line-height: 1.2;
}

.hero p {
    font-size: 18px;
    color: var(--text-secondary);
    max-width: 600px;
    margin: 0 auto;
    line-height: 1.6;
}

/* Status Badge */
.status-badge {
    display: inline-flex;
    align-items: center;
    gap: 8px;
    padding: 6px 14px;
    background: var(--primary-light);
    color: var(--primary-hover);
    border-radius: 100px;
    font-size: 13px;
    font-weight: 500;
    margin-bottom: 20px;
}

.status-dot {
    width: 6px;
    height: 6px;
    background: var(--primary);
    border-radius: 50%;
    animation: pulse 2s infinite;
}

@keyframes pulse {
    0%, 100% { opacity: 1; }
    50% { opacity: 0.5; }
}

/* Main Card */
.card {
    background: var(--surface);
    border: 1px solid var(--border);
    border-radius: var(--radius-lg);
    padding: 32px;
    margin-bottom: 24px;
    transition: var(--transition);
}

.card:hover {
    border-color: #D1D5DB;
    box-shadow: var(--shadow-md);
}

/* Form Elements - Clean Perplexity Style */
.form-group {
    margin-bottom: 24px;
}

.form-row {
    display: grid;
    grid-template-columns: repeat(auto-fit, minmax(200px, 1fr));
    gap: 16px;
    margin-bottom: 24px;
}

label {
    display: block;
    font-size: 14px;
    font-weight: 500;
    color: var(--text);
    margin-bottom: 8px;
}

.input {
    width: 100%;
    padding: 12px 16px;
    background: var(--bg);
    border: 1px solid var(--border);
    border-radius: var(--radius-md);
    color: var(--text);
    font-size: 15px;
    font-family: inherit;
    transition: var(--transition);
    outline: none;
}

.input:hover {
    border-color: #D1D5DB;
}

.input:focus {
    border-color: var(--border-focus);
    box-shadow: 0 0 0 3px rgba(16, 185, 129, 0.1);
}

select.input {
    appearance: none;
    background-image: url("data:image/svg+xml,%3Csvg width='10' height='6' viewBox='0 0 10 6' fill='none' xmlns='http://www.w3.org/2000/svg'%3E%3Cpath d='M1 1L5 5L9 1' stroke='%236B7280' stroke-width='1.5' stroke-linecap='round' stroke-linejoin='round'/%3E%3C/svg%3E");
    background-repeat: no-repeat;
    background-position: right 14px center;
    padding-right: 40px;
    cursor: pointer;
}

textarea.input {
    min-height: 140px;
    resize: vertical;
    font-family: 'SF Mono', Monaco, monospace;
    font-size: 14px;
    line-height: 1.6;
}

/* Upload Zone - Clean Design */
.upload-zone {
    border: 2px dashed var(--border);
    border-radius: var(--radius-md);
    padding: 40px 32px;
    text-align: center;
    cursor: pointer;
    transition: var(--transition);
    background: var(--bg);
    position: relative;
}

.upload-zone:hover, .upload-zone.drag-active {
    border-color: var(--primary);
    background: var(--primary-light);
}

.upload-zone input {
    position: absolute;
    inset: 0;
    opacity: 0;
    cursor: pointer;
    width: 100%;
    height: 100%;
}

.upload-icon {
    width: 48px;
    height: 48px;
    margin: 0 auto 16px;
    background: var(--surface);
    border: 1px solid var(--border);
    border-radius: var(--radius-md);
    display: flex;
    align-items: center;
    justify-content: center;
    font-size: 24px;
    transition: var(--transition);
}

.upload-zone:hover .upload-icon {
    background: var(--primary-light);
    border-color: var(--primary);
    transform: scale(1.05);
}

.upload-zone h4 {
    font-size: 16px;
    font-weight: 500;
    color: var(--text);
    margin-bottom: 6px;
}

.upload-zone p {
    font-size: 14px;
    color: var(--text-muted);
}

.file-preview {
    margin-top: 16px;
    padding: 12px 16px;
    background: var(--primary-light);
    border: 1px solid var(--primary);
    border-radius: var(--radius-md);
    display: none;
    align-items: center;
    gap: 10px;
    font-size: 14px;
    color: var(--primary-hover);
    font-weight: 500;
}

.file-preview.show {
    display: flex;
}

.file-preview button {
    margin-left: auto;
    background: none;
    border: none;
    color: var(--text-muted);
    cursor: pointer;
    font-size: 18px;
    line-height: 1;
    padding: 0;
    width: 24px;
    height: 24px;
    display: flex;
    align-items: center;
    justify-content: center;
    border-radius: 4px;
    transition: var(--transition);
}

.file-preview button:hover {
    background: rgba(0,0,0,0.05);
    color: var(--error);
}

/* Or Divider */
.or-divider {
    display: flex;
    align-items: center;
    gap: 16px;
    margin: 20px 0;
    color: var(--text-muted);
    font-size: 13px;
    font-weight: 500;
    text-transform: uppercase;
}

.or-divider::before, .or-divider::after {
    content: '';
    flex: 1;
    height: 1px;
    background: var(--border);
}

/* Button - Perplexity Style */
.btn {
    display: inline-flex;
    align-items: center;
    justify-content: center;
    gap: 8px;
    padding: 14px 28px;
    background: var(--text);
    color: white;
    border: none;
    border-radius: var(--radius-md);
    font-size: 15px;
    font-weight: 600;
    cursor: pointer;
    transition: var(--transition);
    width: 100%;
}

.btn:hover {
    background: #374151;
    transform: translateY(-1px);
    box-shadow: var(--shadow-md);
}

.btn:active {
    transform: translateY(0);
}

.btn:disabled {
    opacity: 0.6;
    cursor: not-allowed;
    transform: none;
}

.btn-primary {
    background: var(--primary);
}

.btn-primary:hover {
    background: var(--primary-hover);
}

.btn-secondary {
    background: var(--bg);
    color: var(--text);
    border: 1px solid var(--border);
}

.btn-secondary:hover {
    background: var(--surface);
}

.btn-sm {
    padding: 8px 16px;
    font-size: 14px;
    width: auto;
}

/* Loading State - Minimal */
.loading-state {
    display: none;
    text-align: center;
    padding: 60px 20px;
}

.loading-spinner {
    width: 40px;
    height: 40px;
    margin: 0 auto 20px;
    border: 3px solid var(--border);
    border-top-color: var(--primary);
    border-radius: 50%;
    animation: spin 0.8s linear infinite;
}

@keyframes spin {
    to { transform: rotate(360deg); }
}

.loading-state h3 {
    font-size: 18px;
    font-weight: 600;
    color: var(--text);
    margin-bottom: 8px;
}

.loading-state p {
    font-size: 14px;
    color: var(--text-muted);
}

/* Result State */
.result-state {
    display: none;
}

.result-header {
    display: flex;
    align-items: center;
    justify-content: space-between;
    margin-bottom: 20px;
    flex-wrap: wrap;
    gap: 12px;
}

.result-title {
    display: flex;
    align-items: center;
    gap: 12px;
    font-size: 18px;
    font-weight: 600;
}

.success-badge {
    width: 28px;
    height: 28px;
    background: var(--primary);
    color: white;
    border-radius: 50%;
    display: flex;
    align-items: center;
    justify-content: center;
    font-size: 14px;
}

/* Output Box - Clean Code Style */
.output-box {
    background: #1F2937;
    border-radius: var(--radius-md);
    padding: 20px;
    color: #E5E7EB;
    font-family: 'SF Mono', Monaco, monospace;
    font-size: 14px;
    line-height: 1.7;
    max-height: 500px;
    overflow-y: auto;
    white-space: pre-wrap;
    position: relative;
}

.output-box::before {
    content: '';
    position: absolute;
    top: 0;
    left: 0;
    right: 0;
    height: 40px;
    background: linear-gradient(to bottom, rgba(31,41,55,0.8), transparent);
    border-radius: var(--radius-md) var(--radius-md) 0 0;
    pointer-events: none;
}

.output-header {
    display: flex;
    align-items: center;
    gap: 8px;
    margin-bottom: 16px;
    padding-bottom: 12px;
    border-bottom: 1px solid rgba(255,255,255,0.1);
}

.dot {
    width: 12px;
    height: 12px;
    border-radius: 50%;
}

.dot-red { background: #EF4444; }
.dot-yellow { background: #F59E0B; }
.dot-green { background: #10B981; }

/* Features Grid */
.features {
    display: grid;
    grid-template-columns: repeat(auto-fit, minmax(250px, 1fr));
    gap: 20px;
    margin-top: 40px;
}

.feature {
    padding: 24px;
    background: var(--surface);
    border: 1px solid var(--border);
    border-radius: var(--radius-lg);
    transition: var(--transition);
}

.feature:hover {
    border-color: var(--primary);
    transform: translateY(-2px);
    box-shadow: var(--shadow-md);
}

.feature-icon {
    width: 40px;
    height: 40px;
    background: var(--primary-light);
    border-radius: var(--radius-md);
    display: flex;
    align-items: center;
    justify-content: center;
    font-size: 20px;
    margin-bottom: 16px;
}

.feature h4 {
    font-size: 16px;
    font-weight: 600;
    margin-bottom: 6px;
    color: var(--text);
}

.feature p {
    font-size: 14px;
    color: var(--text-secondary);
    line-height: 1.5;
}

/* Footer */
.footer {
    text-align: center;
    padding: 40px 24px;
    color: var(--text-muted);
    font-size: 13px;
    border-top: 1px solid var(--border);
    margin-top: 60px;
}

/* Auth Pages - Clean Centered */
.auth-container {
    min-height: 100vh;
    display: flex;
    align-items: center;
    justify-content: center;
    padding: 24px;
    background: var(--surface);
}

.auth-card {
    width: 100%;
    max-width: 420px;
    background: var(--bg);
    border: 1px solid var(--border);
    border-radius: var(--radius-lg);
    padding: 40px;
    box-shadow: var(--shadow-lg);
}

.auth-header {
    text-align: center;
    margin-bottom: 32px;
}

.auth-logo {
    width: 56px;
    height: 56px;
    margin: 0 auto 20px;
    background: linear-gradient(135deg, var(--primary), #34D399);
    border-radius: var(--radius-md);
    display: flex;
    align-items: center;
    justify-content: center;
    color: white;
    font-size: 24px;
    font-weight: 700;
}

.auth-header h1 {
    font-size: 24px;
    font-weight: 700;
    margin-bottom: 8px;
}

.auth-header p {
    color: var(--text-secondary);
    font-size: 15px;
}

/* Alert */
.alert {
    padding: 12px 16px;
    border-radius: var(--radius-md);
    font-size: 14px;
    margin-bottom: 20px;
    display: flex;
    align-items: center;
    gap: 10px;
}

.alert-error {
    background: var(--error-bg);
    color: var(--error);
    border: 1px solid #FECACA;
}

/* Responsive */
@media (max-width: 640px) {
    .container {
        padding: 80px 16px 32px;
    }
    
    .card {
        padding: 24px;
    }
    
    .form-row {
        grid-template-columns: 1fr;
    }
    
    .nav {
        padding: 0 16px;
    }
    
    .hero h1 {
        font-size: 28px;
    }
}

/* Animations */
@keyframes fadeIn {
    from { opacity: 0; transform: translateY(10px); }
    to { opacity: 1; transform: translateY(0); }
}

.animate-in {
    animation: fadeIn 0.4s ease-out forwards;
}
"""


# ---------------- HTML TEMPLATES ----------------

def get_main_page():
    """Return Perplexity-inspired main page"""
    return f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0, viewport-fit=cover">
    <meta name="theme-color" content="#FFFFFF">
    <title>AI Business Companion</title>
    <style>{PERPLEXITY_CSS}</style>
</head>
<body>
    <nav class="nav">
        <a href="/" class="nav-brand">
            <div class="logo-icon">AI</div>
            <span>Business Companion</span>
        </a>
        <div class="nav-actions">
            <a href="/logout" class="btn btn-secondary btn-sm">Logout</a>
        </div>
    </nav>
    
    <div class="container">
        <div class="hero animate-in">
            <div class="status-badge">
                <span class="status-dot"></span>
                System Operational
            </div>
            <h1>Transform Data Into Decisions</h1>
            <p>Upload any file or paste data to generate comprehensive business intelligence reports instantly.</p>
        </div>
        
        <div class="card animate-in">
            <div class="form-row">
                <div class="form-group" style="margin-bottom:0;">
                    <label>Business Sector</label>
                    <select id="business" class="input">
                        <option>General Business</option>
                        <option>Retail & E-Commerce</option>
                        <option>Real Estate</option>
                        <option>Financial Services</option>
                        <option>Healthcare</option>
                        <option>Technology</option>
                        <option>Manufacturing</option>
                    </select>
                </div>
                
                <div class="form-group" style="margin-bottom:0;">
                    <label>Analysis Focus</label>
                    <select id="focus" class="input">
                        <option value="full">Comprehensive Analysis</option>
                        <option value="profit">Profit Optimization</option>
                        <option value="loss">Risk Assessment</option>
                        <option value="growth">Growth Strategy</option>
                    </select>
                </div>
                
                <div class="form-group" style="margin-bottom:0;">
                    <label>Output Format</label>
                    <select id="outputType" class="input">
                        <option value="text">Interactive Report</option>
                        <option value="pdf">PDF Document</option>
                        <option value="docx">Word Document</option>
                        <option value="pptx">PowerPoint</option>
                    </select>
                </div>
            </div>
            
            <div class="form-group">
                <label>Upload Data</label>
                <div class="upload-zone" id="uploadZone">
                    <input type="file" id="fileInput" accept=".csv,.xlsx,.xls,.json,.txt,.md,.pdf,.docx,.doc">
                    <div class="upload-icon">📄</div>
                    <h4>Drop file here or click to browse</h4>
                    <p>Supports CSV, Excel, PDF, Word, JSON, and text files</p>
                </div>
                <div class="file-preview" id="filePreview">
                    <span>📎</span>
                    <span id="fileName"></span>
                    <button onclick="clearFile()" title="Remove file">×</button>
                </div>
            </div>
            
            <div class="or-divider">or paste data</div>
            
            <div class="form-group" style="margin-bottom:0;">
                <textarea id="dataInput" class="input" placeholder="Paste your data here... CSV, JSON, or plain text"></textarea>
            </div>
            
            <div style="margin-top:24px;">
                <button class="btn btn-primary" id="genBtn" onclick="generate()">
                    <span>⚡</span>
                    <span>Generate Report</span>
                </button>
            </div>
        </div>
        
        <div class="loading-state" id="loadingState">
            <div class="loading-spinner"></div>
            <h3>Processing your data</h3>
            <p id="loadingText">Analyzing patterns...</p>
        </div>
        
        <div class="result-state" id="resultState">
            <div class="card">
                <div class="result-header">
                    <div class="result-title">
                        <div class="success-badge">✓</div>
                        <span>Analysis Complete</span>
                    </div>
                    <div style="display:flex;gap:10px;">
                        <button class="btn btn-secondary btn-sm" onclick="resetForm()">New Analysis</button>
                        <button class="btn btn-primary btn-sm" id="dlBtn" onclick="downloadFile()" style="display:none;">Download</button>
                    </div>
                </div>
                <div class="output-box" id="outputBox">
                    <div class="output-header">
                        <span class="dot dot-red"></span>
                        <span class="dot dot-yellow"></span>
                        <span class="dot dot-green"></span>
                        <span style="margin-left:8px;font-size:12px;color:#9CA3AF;">output</span>
                    </div>
                    <div id="outputContent"></div>
                </div>
            </div>
        </div>
        
        <div class="features">
            <div class="feature">
                <div class="feature-icon">🎯</div>
                <h4>Universal File Support</h4>
                <p>Automatically detects and processes Excel, PDF, Word, JSON, CSV, and text files.</p>
            </div>
            <div class="feature">
                <div class="feature-icon">⚡</div>
                <h4>Instant Analysis</h4>
                <p>AI-powered insights generated in seconds with executive-ready formatting.</p>
            </div>
            <div class="feature">
                <div class="feature-icon">🔒</div>
                <h4>Secure Processing</h4>
                <p>Your data is processed in-memory and never stored on our servers.</p>
            </div>
        </div>
    </div>
    
    <div class="footer">
        <p>AI Business Companion • Professional Intelligence Platform</p>
    </div>
    
    <script>
        const fileInput = document.getElementById('fileInput');
        const uploadZone = document.getElementById('uploadZone');
        const filePreview = document.getElementById('filePreview');
        const fileName = document.getElementById('fileName');
        const dataInput = document.getElementById('dataInput');
        
        // File handling
        uploadZone.addEventListener('dragover', (e) => {{
            e.preventDefault();
            uploadZone.classList.add('drag-active');
        }});
        
        uploadZone.addEventListener('dragleave', () => {{
            uploadZone.classList.remove('drag-active');
        }});
        
        uploadZone.addEventListener('drop', (e) => {{
            e.preventDefault();
            uploadZone.classList.remove('drag-active');
            if (e.dataTransfer.files.length) {{
                handleFile(e.dataTransfer.files[0]);
            }}
        }});
        
        fileInput.addEventListener('change', () => {{
            if (fileInput.files.length) {{
                handleFile(fileInput.files[0]);
            }}
        }});
        
        function handleFile(file) {{
            // Check file size (50MB limit)
            if (file.size > 50 * 1024 * 1024) {{
                alert('File too large. Maximum size is 50MB.');
                return;
            }}
            
            fileName.textContent = file.name + ' (' + formatBytes(file.size) + ')';
            filePreview.classList.add('show');
            dataInput.value = ''; // Clear text input when file selected
        }}
        
        function formatBytes(bytes) {{
            if (bytes === 0) return '0 Bytes';
            const k = 1024;
            const sizes = ['Bytes', 'KB', 'MB', 'GB'];
            const i = Math.floor(Math.log(bytes) / Math.log(k));
            return parseFloat((bytes / Math.pow(k, i)).toFixed(1)) + ' ' + sizes[i];
        }}
        
        function clearFile() {{
            fileInput.value = '';
            filePreview.classList.remove('show');
        }}
        
        // Generate report
        const messages = [
            "Detecting file format...",
            "Extracting and structuring data...",
            "Running AI analysis models...",
            "Generating insights and recommendations...",
            "Finalizing report..."
        ];
        
        async function generate() {{
            const btn = document.getElementById('genBtn');
            const loading = document.getElementById('loadingState');
            const result = document.getElementById('resultState');
            const outputType = document.getElementById('outputType').value;
            
            // Validation
            const hasFile = fileInput.files.length > 0;
            const hasText = dataInput.value.trim().length > 0;
            
            if (!hasFile && !hasText) {{
                uploadZone.style.borderColor = 'var(--error)';
                setTimeout(() => {{
                    uploadZone.style.borderColor = '';
                }}, 2000);
                return;
            }}
            
            // UI State
            btn.disabled = true;
            loading.style.display = 'block';
            result.style.display = 'none';
            loading.scrollIntoView({{behavior: 'smooth', block: 'nearest'}});
            
            // Cycle messages
            let msgIdx = 0;
            const msgInterval = setInterval(() => {{
                document.getElementById('loadingText').textContent = messages[msgIdx];
                msgIdx = (msgIdx + 1) % messages.length;
            }}, 1500);
            
            // Prepare form data
            const formData = new FormData();
            formData.append('business', document.getElementById('business').value);
            formData.append('focus', document.getElementById('focus').value);
            formData.append('output_type', outputType);
            
            if (hasFile) {{
                formData.append('file', fileInput.files[0]);
            }} else {{
                formData.append('text_data', dataInput.value);
            }}
            
            try {{
                const response = await fetch('/generate', {{
                    method: 'POST',
                    body: formData
                }});
                
                clearInterval(msgInterval);
                
                if (outputType === 'text') {{
                    const data = await response.json();
                    if (data.success) {{
                        showResult(data.report, false);
                    }} else {{
                        throw new Error(data.error);
                    }}
                }} else {{
                    const blob = await response.blob();
                    window.currentBlob = blob;
                    window.currentExt = outputType;
                    showResult('✓ Document generated successfully. Click download to save.', true);
                }}
            }} catch (error) {{
                clearInterval(msgInterval);
                showResult('Error: ' + error.message, false);
            }}
            
            btn.disabled = false;
        }}
        
        function showResult(text, isDownloadable) {{
            document.getElementById('loadingState').style.display = 'none';
            document.getElementById('resultState').style.display = 'block';
            document.getElementById('dlBtn').style.display = isDownloadable ? 'inline-flex' : 'none';
            
            const outputContent = document.getElementById('outputContent');
            
            if (isDownloadable) {{
                outputContent.textContent = text;
            }} else {{
                // Typewriter effect for text output
                outputContent.innerHTML = '';
                let i = 0;
                const speed = 5;
                function type() {{
                    if (i < text.length) {{
                        const char = text.charAt(i) === '\\n' ? '<br>' : text.charAt(i);
                        outputContent.innerHTML += char;
                        i++;
                        document.getElementById('outputBox').scrollTop = document.getElementById('outputBox').scrollHeight;
                        setTimeout(type, speed);
                    }}
                }}
                type();
            }}
            
            document.getElementById('resultState').scrollIntoView({{behavior: 'smooth'}});
        }}
        
        function downloadFile() {{
            const url = URL.createObjectURL(window.currentBlob);
            const a = document.createElement('a');
            a.href = url;
            a.download = `Report_${{new Date().toISOString().slice(0,10)}}.${{window.currentExt}}`;
            document.body.appendChild(a);
            a.click();
            a.remove();
            URL.revokeObjectURL(url);
        }}
        
        function resetForm() {{
            fileInput.value = '';
            filePreview.classList.remove('show');
            dataInput.value = '';
            document.getElementById('resultState').style.display = 'none';
            window.scrollTo({{top: 0, behavior: 'smooth'}});
        }}
    </script>
</body>
</html>"""


def get_login_page(error=False):
    """Return clean login page"""
    error_html = '<div class="alert alert-error"><span>⚠️</span> Invalid email or password</div>' if error else ''
    
    return f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Sign In - AI Business Companion</title>
    <style>{PERPLEXITY_CSS}</style>
</head>
<body>
    <div class="auth-container">
        <div class="auth-card">
            <div class="auth-header">
                <div class="auth-logo">AI</div>
                <h1>Welcome back</h1>
                <p>Sign in to access your dashboard</p>
            </div>
            {error_html}
            <form method="POST" action="/login">
                <div class="form-group">
                    <label>Email address</label>
                    <input type="email" name="email" class="input" placeholder="you@company.com" required autofocus>
                </div>
                <div class="form-group">
                    <label>Password</label>
                    <input type="password" name="password" class="input" placeholder="••••••••" required>
                </div>
                <button type="submit" class="btn btn-primary" style="margin-top:8px;">Sign In</button>
            </form>
            <p style="text-align:center;margin-top:24px;font-size:14px;color:var(--text-secondary);">
                Don't have an account? <a href="/signup" style="color:var(--text);font-weight:600;text-decoration:none;">Sign up</a>
            </p>
        </div>
    </div>
</body>
</html>"""


def get_signup_page():
    """Return clean signup page"""
    return f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Create Account - AI Business Companion</title>
    <style>{PERPLEXITY_CSS}</style>
</head>
<body>
    <div class="auth-container">
        <div class="auth-card">
            <div class="auth-header">
                <div class="auth-logo">+</div>
                <h1>Create your account</h1>
                <p>Start generating business intelligence</p>
            </div>
            <form method="POST" action="/signup">
                <div class="form-group">
                    <label>Email address</label>
                    <input type="email" name="email" class="input" placeholder="you@company.com" required autofocus>
                </div>
                <div class="form-group">
                    <label>Password</label>
                    <input type="password" name="password" class="input" placeholder="Min 6 characters" required minlength="6">
                </div>
                <button type="submit" class="btn btn-primary" style="margin-top:8px;">Create Account</button>
            </form>
            <p style="text-align:center;margin-top:24px;font-size:14px;color:var(--text-secondary);">
                Already have an account? <a href="/login" style="color:var(--text);font-weight:600;text-decoration:none;">Sign in</a>
            </p>
        </div>
    </div>
</body>
</html>"""


# ---------------- ROUTES ----------------

@app.route("/")
def index():
    if "user" not in session:
        return redirect("/login")
    return get_main_page()


@app.route("/generate", methods=["POST"])
def generate():
    """Intelligent report generation with universal file support"""
    try:
        output_type = request.form.get("output_type", "text")
        uploaded_file = request.files.get("file")
        text_data = request.form.get("text_data", "").strip()
        
        # Initialize processor
        processor = DataProcessor()
        content = None
        data_type = "text"
        source_info = ""
        
        # Process uploaded file
        if uploaded_file:
            filename = uploaded_file.filename
            raw_data = uploaded_file.read()
            
            if not raw_data:
                return jsonify({"success": False, "error": "File is empty"}), 400
            
            # Check file size
            if len(raw_data) > 50 * 1024 * 1024:
                return jsonify({"success": False, "error": "File too large (max 50MB)"}), 400
            
            # Detect and extract
            try:
                detected_format = processor.detect_format(raw_data, filename)
                content = processor.extract_text(raw_data, filename, detected_format)
                data_type = detected_format if detected_format not in ["binary", "empty"] else "text"
                source_info = f"File: {filename}"
            except ValueError as e:
                return jsonify({"success": False, "error": str(e)}), 400
            except Exception as e:
                app.logger.error(f"File processing error: {str(e)}")
                return jsonify({"success": False, "error": f"Could not process file: {str(e)}"}), 400
        
        # Process pasted text
        elif text_data:
            content = text_data
            detected = processor.detect_format(content.encode(), "")
            data_type = detected if detected != "binary" else "text"
            source_info = "Pasted text"
        else:
            return jsonify({"success": False, "error": "Please provide a file or paste data"}), 400
        
        # Validate content
        if len(content) > 10_000_000:  # 10MB text limit
            content = content[:10_000_000] + "\n\n[Content truncated due to size]"
        
        if not content.strip():
            return jsonify({"success": False, "error": "No readable content found in input"}), 400
        
        # Generate report
        report = generator.generate_report(
            data_content=content,
            data_type=data_type,
            user_name=session.get("user", "User"),
            user_role="Business Analyst",
            business_type=request.form.get("business", "General"),
            report_focus=request.form.get("focus", "full")
        )
        
        # Add source attribution
        report = f"Source: {source_info}\nFormat: {data_type.upper()}\n\n{report}"
        
        # Generate output
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        
        if output_type == "pdf":
            # Create in temp directory
            temp_dir = tempfile.mkdtemp()
            filename = os.path.join(temp_dir, f"report_{timestamp}.pdf")
            try:
                create_pdf(report, filename)
                return send_file(filename, 
                               as_attachment=True, 
                               download_name=f"Business_Report_{timestamp}.pdf",
                               mimetype='application/pdf')
            finally:
                shutil.rmtree(temp_dir, ignore_errors=True)
        
        elif output_type == "docx":
            temp_dir = tempfile.mkdtemp()
            filename = os.path.join(temp_dir, f"report_{timestamp}.docx")
            try:
                create_docx(report, filename)
                return send_file(filename,
                               as_attachment=True,
                               download_name=f"Business_Report_{timestamp}.docx",
                               mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document')
            finally:
                shutil.rmtree(temp_dir, ignore_errors=True)
        
        elif output_type == "pptx":
            temp_dir = tempfile.mkdtemp()
            filename = os.path.join(temp_dir, f"report_{timestamp}.pptx")
            try:
                create_ppt(report, filename)
                return send_file(filename,
                               as_attachment=True,
                               download_name=f"Business_Report_{timestamp}.pptx",
                               mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')
            finally:
                shutil.rmtree(temp_dir, ignore_errors=True)
        
        else:  # text
            return jsonify({
                "success": True,
                "report": report,
                "meta": {
                    "source": source_info,
                    "format": data_type,
                    "timestamp": timestamp
                }
            })
            
    except Exception as e:
        app.logger.error(f"Generation error: {str(e)}")
        return jsonify({"success": False, "error": f"Processing failed: {str(e)}"}), 500


@app.route("/login", methods=["GET", "POST"])
def login():
    if request.method == "POST":
        email = request.form.get("email", "").strip().lower()
        password = request.form.get("password", "")
        
        if not email or not password:
            return get_login_page(error=True)
        
        if validate_user(email, password):
            session.permanent = True
            session["user"] = email
            return redirect("/")
        
        return get_login_page(error=True)
    
    return get_login_page(error=False)


@app.route("/signup", methods=["GET", "POST"])
def signup():
    if request.method == "POST":
        email = request.form.get("email", "").strip().lower()
        password = request.form.get("password", "")
        
        if not email or not password:
            return "Email and password required", 400
        
        if len(password) < 6:
            return get_signup_page()  # Could add error message here
        
        try:
            create_user(email, password)
            return redirect("/login")
        except Exception as e:
            # Check if user exists
            if "already" in str(e).lower() or "exists" in str(e).lower():
                return get_login_page(error=False)  # Redirect to login
            return f"Error: {str(e)}", 400
    
    return get_signup_page()


@app.route("/logout")
def logout():
    session.clear()
    return redirect("/login")


@app.route("/health")
def health_check():
    """Health check endpoint"""
    return jsonify({
        "status": "healthy",
        "version": "2.0",
        "features": ["universal_file_support", "auto_format_detection", "multi_format_export"]
    })


if __name__ == "__main__":
    print("""
    ╔══════════════════════════════════════════════════════════════╗
    ║         AI Business Companion - Professional Edition         ║
    ║                                                              ║
    ║  ✓ Universal file support (Excel, PDF, Word, JSON, CSV)     ║
    ║  ✓ Auto-format detection                                     ║
    ║  ✓ Perplexity AI-inspired minimal design                     ║
    ║  ✓ Optimized for mobile, tablet, and desktop                ║
    ║                                                              ║
    ║  Running at http://localhost:5000                            ║
    ╚══════════════════════════════════════════════════════════════╝
    """)
    app.run(debug=True, host='0.0.0.0', port=5000, threaded=True)